
from __future__ import annotations

from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple
import json
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400


@dataclass
class PptBlock:
    block_id: str
    source_file: str
    source_path: str
    doc_type: str
    block_type: str
    order: int
    text: str
    section_path: List[str] = field(default_factory=list)
    page_number: Optional[int] = None
    slide_number: Optional[int] = None
    shape_index: Optional[int] = None
    sheet_name: Optional[str] = None
    row_number: Optional[int] = None
    metadata: Dict[str, Any] = field(default_factory=dict)


def build_section_path(*parts: Optional[str]) -> List[str]:
    return [part for part in parts if part]


class PptParser:
    """
    PPTX parser focused on RAG-friendly extraction.

    Handles:
    - regular text boxes
    - title placeholders
    - bullet paragraphs
    - tables
    - grouped shapes (recursive flattening)
    - notes (optional)
    - image placeholders (optional)

    Important:
    PowerPoint does not guarantee a true reading order, so this parser
    approximates order using visual position: top -> left.
    """

    def __init__(
        self,
        include_notes: bool = True,
        include_images: bool = True,
        include_empty_tables: bool = False,
    ) -> None:
        self.include_notes = include_notes
        self.include_images = include_images
        self.include_empty_tables = include_empty_tables

    def parse(self, file_path: str) -> List[PptBlock]:
        prs = Presentation(file_path)
        source_file = Path(file_path).name
        source_path = str(Path(file_path).resolve())
        blocks: List[PptBlock] = []
        order = 0

        for slide_number, slide in enumerate(prs.slides, start=1):
            section_hint = self._get_slide_title(slide)

            atomic_shapes = self._flatten_shapes(slide.shapes)
            atomic_shapes.sort(key=self._visual_sort_key)

            for shape_index, item in enumerate(atomic_shapes, start=1):
                shape = item["shape"]
                parent_group = item["parent_group"]
                shape_path = item["shape_path"]

                new_blocks = self._parse_atomic_shape(
                    shape=shape,
                    source_file=source_file,
                    source_path=source_path,
                    slide_number=slide_number,
                    shape_index=shape_index,
                    parent_group=parent_group,
                    shape_path=shape_path,
                    section_hint=section_hint,
                    order_start=order,
                )
                blocks.extend(new_blocks)
                order += len(new_blocks)

            if self.include_notes:
                note_text = self._extract_notes_text(slide)
                if note_text:
                    order += 1
                    blocks.append(
                        PptBlock(
                            block_id=f"{source_file}::slide_{slide_number}::notes",
                            source_file=source_file,
                            source_path=source_path,
                            doc_type="pptx",
                            block_type="notes",
                            order=order,
                            text=note_text,
                            section_path=build_section_path(section_hint),
                            slide_number=slide_number,
                            shape_index=9999,
                            metadata={
                                "source": "notes_slide",
                                "parent_group": None,
                            },
                        )
                    )

        return blocks

    def _flatten_shapes(
        self,
        shapes: Iterable[Any],
        parent_group: Optional[str] = None,
        shape_path: str = "",
    ) -> List[Dict[str, Any]]:
        flat: List[Dict[str, Any]] = []

        for shape in shapes:
            current_path = f"{shape_path}/{getattr(shape, 'name', 'shape')}"
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                flat.extend(
                    self._flatten_shapes(
                        shape.shapes,
                        parent_group=getattr(shape, "name", None),
                        shape_path=current_path,
                    )
                )
            else:
                flat.append(
                    {
                        "shape": shape,
                        "parent_group": parent_group,
                        "shape_path": current_path,
                    }
                )

        return flat

    def _parse_atomic_shape(
        self,
        shape: Any,
        source_file: str,
        source_path: str,
        slide_number: int,
        shape_index: int,
        parent_group: Optional[str],
        shape_path: str,
        section_hint: Optional[str],
        order_start: int,
    ) -> List[PptBlock]:
        blocks: List[PptBlock] = []

        if self._should_skip_shape(shape):
            return blocks

        metadata = self._shape_metadata(shape, shape_path=shape_path)

        if shape.has_table:
            table_text = self._extract_table_text(shape)
            if table_text or self.include_empty_tables:
                blocks.append(
                    PptBlock(
                        block_id=f"{source_file}::slide_{slide_number}::shape_{shape_index}",
                        source_file=source_file,
                        source_path=source_path,
                        doc_type="pptx",
                        block_type="table",
                        order=order_start + len(blocks) + 1,
                        text=table_text,
                        section_path=build_section_path(section_hint),
                        slide_number=slide_number,
                        shape_index=shape_index,
                        metadata={
                            **metadata,
                            "parent_group": parent_group,
                        },
                    )
                )
            return blocks

        if getattr(shape, "has_text_frame", False):
            blocks.extend(
                self._extract_text_blocks(
                    shape=shape,
                    source_file=source_file,
                    source_path=source_path,
                    slide_number=slide_number,
                    shape_index=shape_index,
                    parent_group=parent_group,
                    metadata=metadata,
                    section_hint=section_hint,
                    order_start=order_start + len(blocks),
                )
            )

        elif self.include_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blocks.append(
                PptBlock(
                    block_id=f"{source_file}::slide_{slide_number}::shape_{shape_index}",
                    source_file=source_file,
                    source_path=source_path,
                    doc_type="pptx",
                    block_type="image",
                    order=order_start + len(blocks) + 1,
                    text="[image]",
                    section_path=build_section_path(section_hint),
                    slide_number=slide_number,
                    shape_index=shape_index,
                    metadata={
                        **metadata,
                        "parent_group": parent_group,
                    },
                )
            )

        return blocks

    def _extract_text_blocks(
        self,
        shape: Any,
        source_file: str,
        source_path: str,
        slide_number: int,
        shape_index: int,
        parent_group: Optional[str],
        metadata: Dict[str, Any],
        section_hint: Optional[str],
        order_start: int,
    ) -> List[PptBlock]:
        paragraphs = list(getattr(shape.text_frame, "paragraphs", []))
        if not paragraphs:
            return []

        title_lines: List[str] = []
        bullet_lines: List[str] = []
        body_lines: List[str] = []

        for para in paragraphs:
            raw_text = "".join(run.text for run in para.runs) if para.runs else para.text
            text = self._normalize_text(raw_text)
            if not text:
                continue

            para_level = getattr(para, "level", 0) or 0

            if self._looks_like_title(shape):
                title_lines.append(text)
            elif para_level > 0 or self._looks_like_bullet_text(text):
                bullet_lines.append(("  " * para_level) + text)
            else:
                body_lines.append(text)

        blocks: List[PptBlock] = []
        base_id = f"{source_file}::slide_{slide_number}::shape_{shape_index}"
        section_path = build_section_path(section_hint)
        normalized_metadata = {
            **metadata,
            "parent_group": parent_group,
        }

        if title_lines:
            blocks.append(
                PptBlock(
                    block_id=base_id + "::title",
                    source_file=source_file,
                    source_path=source_path,
                    doc_type="pptx",
                    block_type="title",
                    order=order_start + len(blocks) + 1,
                    text="\n".join(title_lines),
                    section_path=section_path,
                    slide_number=slide_number,
                    shape_index=shape_index,
                    metadata=normalized_metadata,
                )
            )

        if body_lines:
            block_type = self._classify_body_text(shape)
            blocks.append(
                PptBlock(
                    block_id=base_id + f"::{block_type}",
                    source_file=source_file,
                    source_path=source_path,
                    doc_type="pptx",
                    block_type=block_type,
                    order=order_start + len(blocks) + 1,
                    text="\n".join(body_lines),
                    section_path=section_path,
                    slide_number=slide_number,
                    shape_index=shape_index,
                    metadata=normalized_metadata,
                )
            )

        if bullet_lines:
            blocks.append(
                PptBlock(
                    block_id=base_id + "::bullets",
                    source_file=source_file,
                    source_path=source_path,
                    doc_type="pptx",
                    block_type="bullet",
                    order=order_start + len(blocks) + 1,
                    text="\n".join(bullet_lines),
                    section_path=section_path,
                    slide_number=slide_number,
                    shape_index=shape_index,
                    metadata=normalized_metadata,
                )
            )

        return blocks

    def _classify_body_text(self, shape: Any) -> str:
        name = (getattr(shape, "name", "") or "").lower()
        ph_type = self._placeholder_type(shape)

        if "callout" in name:
            return "callout"
        if ph_type in {"SUBTITLE"}:
            return "subtitle"
        return "text"

    def _extract_table_text(self, shape: Any) -> str:
        rows_out: List[str] = []

        for row in shape.table.rows:
            values = [self._normalize_text(cell.text) for cell in row.cells]
            if any(values):
                rows_out.append(" | ".join(values))

        return "\n".join(rows_out).strip()

    def _extract_notes_text(self, slide: Any) -> str:
        try:
            texts = []
            for shape in slide.notes_slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue

                text = self._normalize_text(shape.text)
                if not text:
                    continue

                lower = text.lower()
                if lower in {"click to add notes", "notes"}:
                    continue
                if re.fullmatch(r"\d+", text):
                    continue

                texts.append(text)

            return "\n".join(texts).strip()
        except Exception:
            return ""

    def _get_slide_title(self, slide: Any) -> Optional[str]:
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = self._normalize_text(slide.shapes.title.text)
                return title or None
        except Exception:
            return None
        return None

    def _shape_metadata(self, shape: Any, shape_path: str) -> Dict[str, Any]:
        return {
            "shape_name": getattr(shape, "name", None),
            "shape_path": shape_path,
            "shape_type": str(getattr(shape, "shape_type", None)),
            "placeholder_type": self._placeholder_type(shape),
            "left": int(getattr(shape, "left", 0)),
            "top": int(getattr(shape, "top", 0)),
            "width": int(getattr(shape, "width", 0)),
            "height": int(getattr(shape, "height", 0)),
        }

    def _placeholder_type(self, shape: Any) -> Optional[str]:
        try:
            if getattr(shape, "is_placeholder", False):
                return str(shape.placeholder_format.type).split(" ")[0]
        except Exception:
            return None
        return None

    def _should_skip_shape(self, shape: Any) -> bool:
        ph_type = self._placeholder_type(shape)
        return ph_type in {"DATE", "FOOTER", "SLIDE_NUMBER", "ORG_CHART", "OBJECT"}

    def _looks_like_title(self, shape: Any) -> bool:
        ph_type = self._placeholder_type(shape)
        if ph_type in {"TITLE", "CENTER_TITLE"}:
            return True

        name = (getattr(shape, "name", "") or "").lower()
        if "title" in name:
            return True

        return False

    def _looks_like_bullet_text(self, text: str) -> bool:
        return bool(re.match(r"^([•\-\*\u2013\u2014]|\d+[\.\)])\s+", text))

    def _normalize_text(self, text: Optional[str]) -> str:
        if not text:
            return ""

        text = text.replace("\x0b", " ")
        text = text.replace("\r", "\n")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{2,}", "\n", text)

        return text.strip()

    def _visual_sort_key(self, item: Dict[str, Any]) -> Tuple[int, int, int]:
        shape = item["shape"]
        top = int(getattr(shape, "top", 0))
        left = int(getattr(shape, "left", 0))

        row_bucket = round(top / 50000)
        return (row_bucket, left, top)


def parse_ppt_folder(folder_path: str, parser: Optional[PptParser] = None) -> Dict[str, List[PptBlock]]:
    parser = parser or PptParser()
    folder = Path(folder_path)
    parsed_results: Dict[str, List[PptBlock]] = {}

    for file_path in sorted(folder.glob("*.pptx")):
        if file_path.name.startswith("~$"):
            continue
        parsed_results[file_path.name] = parser.parse(str(file_path))

    return parsed_results


def save_results_to_json(parsed_results: Dict[str, List[PptBlock]], output_file: str = "parsed_ppt_output.json") -> None:
    output_path = Path(output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    json_ready: Dict[str, Any] = {}

    for file_name, blocks in parsed_results.items():
        json_ready[file_name] = [asdict(block) for block in blocks]

    with output_path.open("w", encoding="utf-8") as f:
        json.dump(json_ready, f, indent=2, ensure_ascii=False)


def build_slide_level_text(parsed_results: Dict[str, List[PptBlock]]) -> Dict[str, Any]:
    """
    Builds a second output that combines extracted text at the slide level.
    Useful for quick inspection and early RAG experiments.
    """
    slide_map: Dict[str, Dict[int, List[PptBlock]]] = {}

    for file_name, blocks in parsed_results.items():
        slide_map[file_name] = {}
        for block in blocks:
            slide_map[file_name].setdefault(block.slide_number, []).append(block)

    out: Dict[str, Any] = {}

    for file_name, by_slide in slide_map.items():
        out[file_name] = []
        for slide_number, blocks in sorted(by_slide.items()):
            ordered = sorted(blocks, key=lambda b: (b.shape_index, b.block_id))
            texts = [
                block.text
                for block in ordered
                if block.block_type not in {"image"} and block.text.strip()
            ]

            out[file_name].append(
                {
                    "slide_number": slide_number,
                    "section_hint": ordered[0].section_path[0] if ordered and ordered[0].section_path else None,
                    "text": "\n".join(texts).strip(),
                }
            )

    return out
